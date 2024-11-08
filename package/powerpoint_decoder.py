import pandas as pd
from pptx import Presentation

def list_placeholders_in_layout(slide_layout):
    """
    Prints the index and type of all placeholders in the specified slide layout.
    """
    for placeholder in slide_layout.placeholders:
        print(f"Placeholder index: {placeholder.placeholder_format.idx}, Type: {placeholder.placeholder_format.type}")

def get_placeholder_indexes(slide_layout):
    """
    Extracts and returns placeholder indexes from a given slide layout.
    
    :param slide_layout: A slide layout from which to extract placeholder indexes.
    :return: A list of placeholder indexes.
    """
    return [placeholder.placeholder_format.idx for placeholder in slide_layout.placeholders]


def dataframe_to_slides_content(df, placeholders_per_slide):
    """
    Organizes DataFrame content to fill PowerPoint slides.

    Each slide will be populated with data from the DataFrame columns until all data is used.

    :param df: DataFrame containing the data.
    :param placeholders_per_slide: Number of placeholders available on each slide.
    :return: A list of content for each slide, with each element being a list of content for that slide's placeholders.
    """
    # Flatten the DataFrame into a single list of strings
    data_flat = df.astype(str).values.flatten()

    # Split the flat list into chunks corresponding to the number of placeholders per slide
    content_for_slides = [data_flat[i:i + placeholders_per_slide] for i in range(0, len(data_flat), placeholders_per_slide)]

    return content_for_slides




def populate_template_with_text(template_path, content_data):
    """
    Automatically recognizes placeholder indexes and fills a PowerPoint template with text content.
    
    :param template_path: Path to the PowerPoint template file.
    :param content_data: A list of lists, each sub-list containing text content for placeholders in order.
    """
    # Load the presentation template
    presentation = Presentation(template_path)

    # Choose a specific layout (adjust index as necessary for your template)
    slide_layout = presentation.slide_layouts[0]

    # Extract placeholder indexes from the chosen layout
    placeholder_indexes = get_placeholder_indexes(slide_layout)

    # Iterate through content data to add slides and populate them
    for slide_texts in content_data:
        slide = presentation.slides.add_slide(slide_layout)
        
        for text, placeholder_idx in zip(slide_texts, placeholder_indexes[:len(slide_texts)]):
            placeholder = slide.placeholders[placeholder_idx]
            placeholder.text = text

    # Save the populated presentation
    presentation.save('populated_presentation.pptx')




if __name__ == "__main__":
    template_path = 'Template.pptx'

    # Example DataFrame
    df = pd.DataFrame({
        'Rang': [1, 2, 3, 4, 5],
        'Bib': [342, 238, 563, 322, 561],
        'Name': ['Menk Rene Pascal', 'Donner Christian', 'Pauer Torsten', 'Janas Paul Jan', 'Götze Sebastian'],
        'Verein': ['LAZ Wuppertal', 'LAV Ribnitz-Damgarten/Sanitz', 'Sportclub Magdeburg e.V.', 'Dürener TV 1847', 'Sportclub Magdeburg e.V.'],
        'LV': ['NO', 'MV', 'ST', 'NO', 'ST'],
        'JG': [1983, 1980, 1982, 1980, 1984],
        'Ergebnis': ['7,45', '7,51', '7,75', '7,79', '7,88'],
        'Klasse': ['M40', 'M40', 'M40', 'M40', 'M40'],
        'Info': [' ', 'PB', 'PB', ' ', ' ']
    })
    
    # Calculate the number of placeholders per slide
    placeholders_per_slide = len(get_placeholder_indexes(Presentation(template_path).slide_layouts[0]))
    
    # Prepare the content for slides
    content_data = dataframe_to_slides_content(df, placeholders_per_slide)

    # Populate the template with the content
    populate_template_with_text(template_path, content_data)


